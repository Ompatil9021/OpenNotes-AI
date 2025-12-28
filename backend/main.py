import os
import io
import uvicorn
import uuid
import pandas as pd
import docx
from pptx import Presentation  # <--- NEW: For PowerPoint
from reportlab.pdfgen import canvas # <--- NEW: To save written notes as PDF
from reportlab.lib.pagesizes import letter
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Body
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel  # <--- THIS WAS MISSING
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload
import firebase_admin
from firebase_admin import credentials, firestore
import PyPDF2
from google import genai
from dotenv import load_dotenv

# --- 1. CONFIGURATION (LOAD KEYS FIRST) ---

# Load environment variables from .env file
load_dotenv()

GOOGLE_DRIVE_FOLDER_ID = os.getenv("GOOGLE_DRIVE_FOLDER_ID")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Check if keys exist
if not GEMINI_API_KEY or not GOOGLE_DRIVE_FOLDER_ID:
    raise ValueError("❌ Missing Keys! Make sure the .env file exists and contains GEMINI_API_KEY and GOOGLE_DRIVE_FOLDER_ID.")

# Configure AI Client (After loading key)
ai_client = genai.Client(api_key=GEMINI_API_KEY)

# Service Account Scopes
SCOPES = ['https://www.googleapis.com/auth/drive.file']

# --- 2. SETUP SERVICES ---

if not os.path.exists("serviceAccountKey.json"):
    print("⚠️ WARNING: serviceAccountKey.json not found. Database will fail.")
else:
    try:
        cred = credentials.Certificate("serviceAccountKey.json")
        firebase_admin.initialize_app(cred)
        print("✅ Firebase Database Connected!")
    except ValueError:
        pass 

db = firestore.client()

def get_drive_service():
    if not os.path.exists('token.json'):
        raise Exception("❌ token.json not found. Run setup_drive.py first.")
    creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    return build('drive', 'v3', credentials=creds)
    # Initialize the global variable so the upload function can use it
drive_service = get_drive_service()

# --- 3. HELPER FUNCTIONS ---

def extract_text_from_file(file_bytes, filename):
    try:
        text = ""
        filename = filename.lower()
        file_stream = io.BytesIO(file_bytes)

        if filename.endswith('.pdf'):
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(file_stream)
            for page in pdf_reader.pages: text += page.extract_text() + "\n"

        elif filename.endswith(('.xlsx', '.xls')):
            df_dict = pd.read_excel(file_stream, sheet_name=None)
            for sheet, df in df_dict.items(): text += f"\n--- {sheet} ---\n{df.to_string()}\n"

        elif filename.endswith('.docx'):
            doc = docx.Document(file_stream)
            for para in doc.paragraphs: text += para.text + "\n"

        elif filename.endswith('.pptx'): # 📽️ NEW: PowerPoint Support
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"

        elif filename.endswith(('.txt', '.md')): # 📝 NEW: Markdown/Text
            text = file_bytes.decode('utf-8')

        return text
    except Exception as e:
        print(f"Extraction Error: {e}")
        return ""

def create_pdf_from_text(text, filename):
    """Generates a simple PDF from text string"""
    buffer = io.BytesIO()
    p = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    y = height - 40
    p.setFont("Helvetica", 12)
    
    for line in text.split('\n'):
        if y < 40:
            p.showPage()
            y = height - 40
        # Simple text wrapping could be added here, but for now strict line break
        p.drawString(40, y, line[:90]) # Crop long lines to prevent crash
        y -= 15
        
    p.save()
    buffer.seek(0)
    return buffer
# --- 4. APP API ---
app = FastAPI(title="OpenNotes AI Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- 5. UPDATED UPLOAD ENDPOINT ---
@app.post("/upload")
async def upload_note(
    file: UploadFile = File(...),
    title: str = Form(...),
    subject: str = Form(...),
    course: str = Form(""),       # 🆕
    topic: str = Form(""),        # 🆕
    tags: str = Form(""),         # 🆕
    academic_level: str = Form(""), # 🆕
    description: str = Form(""),
    youtube_url: str = Form(""),
    uploader_id: str = Form(...)
):
    try:
        # 1. Create a safe filename
        safe_filename = f"{uuid.uuid4()}_{file.filename}"
        
        # 2. Read file content
        file_content = await file.read()
        
        # 3. Upload to Google Drive
        # Create a temp file to upload
        with open(safe_filename, "wb") as f:
            f.write(file_content)
            
        g_file = drive_service.files().create(
            body={'name': safe_filename, 'parents': [GOOGLE_DRIVE_FOLDER_ID]},
            # 👇 DYNAMIC MIMETYPE: Allows Drive to preview Excel/Word correctly
            media_body=MediaIoBaseUpload(io.BytesIO(file_content), mimetype=file.content_type),
            fields='id, webViewLink'
        ).execute()
        
        # Clean up temp file
        os.remove(safe_filename)
        
        # 4. Make Public (Viewer)
        drive_service.permissions().create(
            fileId=g_file.get('id'),
            body={'type': 'anyone', 'role': 'reader'}
        ).execute()

        # 👇 CRITICAL STEP: Extract Text for AI (Excel, Word, PDF)
        processed_text = extract_text_from_file(file_content, file.filename)

        # 5. Save to Firestore (MARKED AS PENDING)
        note_data = {
            "title": title,
            "subject": subject,
            "course": course,               # <--- NEW
            "topic": topic,                 # <--- NEW
            "tags": [t.strip() for t in tags.split(',')] if tags else [], # <--- NEW
            "academic_level": academic_level, # <--- NEW
            "description": description or "",
            "youtube_url": youtube_url or "",
            "fileUrl": g_file.get('webViewLink'),
            "driveId": g_file.get('id'),
            "uploader_id": uploader_id,
            "is_approved": False, 
            "processed_text": processed_text,
            "created_at": firestore.SERVER_TIMESTAMP
        }
        
        db.collection("notes").add(note_data)

        return {"status": "success", "message": "Note uploaded successfully! Waiting for approval."}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
# B. CREATE ONLINE NOTE (New Endpoint)
class NoteCreateRequest(BaseModel):
    title: str
    content: str  # The written text
    subject: str
    course: str
    topic: str
    tags: str
    academic_level: str
    uploader_id: str


@app.post("/create-note")
async def create_online_note(req: NoteCreateRequest):
    try:
        # 1. Convert text to PDF so it looks like a "File" in the system
        safe_filename = f"{uuid.uuid4()}_{req.title}.pdf"
        pdf_buffer = create_pdf_from_text(req.content, safe_filename)
        
        # 2. Upload generated PDF to Drive
        media = MediaIoBaseUpload(pdf_buffer, mimetype='application/pdf')
        g_file = drive_service.files().create(
            body={'name': safe_filename, 'parents': [GOOGLE_DRIVE_FOLDER_ID]},
            media_body=media, fields='id, webViewLink'
        ).execute()
        
        drive_service.permissions().create(fileId=g_file.get('id'), body={'type': 'anyone', 'role': 'reader'}).execute()

        # 3. Save to DB (Use the raw content for AI)
        db.collection("notes").add({
            "title": req.title, "subject": req.subject, "course": req.course, "topic": req.topic,
            "tags": [t.strip() for t in req.tags.split(',')] if req.tags else [],
            "academic_level": req.academic_level,
            "description": "Created using Online Editor", 
            "fileUrl": g_file.get('webViewLink'), "driveId": g_file.get('id'),
            "processed_text": req.content, # Direct text!
            "uploader_id": req.uploader_id,
            "is_approved": False, "created_at": firestore.SERVER_TIMESTAMP
        })
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))    
# --- CHAT ENDPOINT ---
@app.post("/chat")
async def chat_with_notes(
    question: str = Body(..., embed=True),
    topic_id: str = Body(..., embed=True)
):
    try:
        print(f"🤖 User asked: {question} (Topic: {topic_id})")

        # 1. Fetch notes
        # NEW WAY: Query the main 'notes' collection by Subject
        # If topic_id is "general" or empty, we fetch everything (optional logic)
        notes_ref = db.collection("notes")
        
        # If the frontend sends a specific subject (e.g., "Computer Science"), filter by it
        if topic_id and topic_id != "general":
            query = notes_ref.where("subject", "==", topic_id)
            docs = query.stream()
        else:
            docs = notes_ref.stream()
        
        # 2. Combine text
        context_text = ""
        for doc in docs:
            data = doc.to_dict()
            if "processed_text" in data:
                context_text += data["processed_text"] + "\n\n"
        
        if not context_text:
            return {"answer": "I couldn't find any notes for this topic. Please upload a PDF first!"}

        # 3. Build Prompt
        prompt = f"""
        You are a helpful tutor. Answer the question based ONLY on the following notes.
        If the answer is not in the notes, say "I don't see that in your notes."
        
        NOTES:
        {context_text[:30000]}
        
        QUESTION:
        {question}
        """

        # 4. Ask Gemini (Using the Free-Tier Friendly Model)
        response = ai_client.models.generate_content(
            model='gemini-flash-latest',
            contents=prompt
        )
        
        print("✅ AI Answered!")
        return {"answer": response.text}

    except Exception as e:
        print(f"❌ Chat Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    
# --- 6. USER DASHBOARD ENDPOINTS ---

# Get notes uploaded by a specific user
@app.get("/my-notes/{user_id}")
async def get_user_notes(user_id: str):
    try:
        # Query Firestore: Get notes where uploader_id matches
        notes_ref = db.collection("notes")
        query = notes_ref.where("uploader_id", "==", user_id)
        docs = query.stream()

        my_notes = []
        for doc in docs:
            note = doc.to_dict()
            note['id'] = doc.id # Add the ID so we can delete it later
            my_notes.append(note)

        return my_notes
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Delete a note
@app.delete("/notes/{note_id}")
async def delete_note(note_id: str):
    try:
        # 1. Delete from Firestore
        db.collection("notes").document(note_id).delete()
        
        # (Optional: In a real app, you would also delete the file from Google Drive here)
        
        return {"status": "success", "message": "Note deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))    

# --- 7. ADMIN ENDPOINTS (FIXED) ---

@app.get("/admin/stats")
async def get_admin_stats():
    try:
        # 1. Fetch the stream ONCE
        notes_stream = db.collection("notes").stream()
        
        # 2. Build the list immediately
        all_notes_data = []
        for doc in notes_stream:
            note = doc.to_dict()
            note['id'] = doc.id  # Attach the ID
            all_notes_data.append(note)
            
        # 3. Return the pre-built list
        return {
            "total_notes": len(all_notes_data),
            "all_notes": all_notes_data
        }
    except Exception as e:
        print(f"Admin Stats Error: {e}") # Print error to console for debugging
        raise HTTPException(status_code=500, detail=str(e))

# --- 8. DYNAMIC SUBJECTS & APPROVALS ---

class SubjectRequest(BaseModel):
    title: str
    code: str
    field: str  # e.g., "Engineering", "Pharmacy"
    description: str
    uploader_id: str

@app.post("/request-subject")
async def request_subject(request: SubjectRequest):
    try:
        # Check if exists (optional logic could go here)
        new_subject = {
            "title": request.title,
            "code": request.code, # <--- 🆕 Save it to database
            "field": request.field,
            "description": request.description,
            "icon": "📚",
            "chapters": [],
            "is_approved": False,
            "uploader_id": request.uploader_id,
            "created_at": firestore.SERVER_TIMESTAMP
        }
        db.collection("subjects").add(new_subject)
        return {"status": "success", "message": "Subject requested! Waiting for Admin approval."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/subjects")
async def get_subjects(show_all: bool = False):
    try:
        subjects_ref = db.collection("subjects")
        
        # If not admin (show_all=False), only show APPROVED subjects
        if not show_all:
            query = subjects_ref.where("is_approved", "==", True)
            docs = query.stream()
        else:
            # Admin sees everything (pending + approved)
            docs = subjects_ref.stream()

        result = []
        for doc in docs:
            sub = doc.to_dict()
            sub['id'] = doc.id
            result.append(sub)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/approve/{collection_name}/{item_id}")
async def approve_item(collection_name: str, item_id: str):
    try:
        # Generic approver for both 'notes' and 'subjects'
        db.collection(collection_name).document(item_id).update({"is_approved": True})
        return {"status": "success", "message": "Item Approved!"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

# --- 9. USER SUBSCRIPTIONS (PERSONALIZATION) ---

class SubscriptionRequest(BaseModel):
    user_id: str
    subject_id: str
    subject_title: str
    subject_desc: str
    subject_icon: str = "📚"

@app.post("/subscribe")
async def subscribe_subject(req: SubscriptionRequest):
    try:
        # Save to: users/{user_id}/subscriptions/{subject_id}
        # This creates a personal list for this user
        user_ref = db.collection("users").document(req.user_id)
        user_ref.collection("subscriptions").document(req.subject_id).set({
            "title": req.subject_title,
            "description": req.subject_desc,
            "icon": req.subject_icon
        })
        return {"status": "success", "message": "Subscribed!"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/subscriptions/{user_id}")
async def get_subscriptions(user_id: str):
    try:
        # Fetch only this user's subjects
        docs = db.collection("users").document(user_id).collection("subscriptions").stream()
        return [{**doc.to_dict(), "id": doc.id} for doc in docs]
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
# --- 10. DELETE SUBJECT (CASCADE DELETE) ---
@app.delete("/subjects/{subject_id}")
async def delete_subject_cascade(subject_id: str):
    try:
        # 1. Get the subject details first (to know its title)
        sub_ref = db.collection("subjects").document(subject_id)
        sub_doc = sub_ref.get()
        
        if not sub_doc.exists:
            raise HTTPException(status_code=404, detail="Subject not found")
            
        subject_title = sub_doc.to_dict().get("title")

        # 2. Find all notes with this subject title
        # Note: We batch delete to be efficient
        notes_query = db.collection("notes").where("subject", "==", subject_title).stream()
        
        deleted_count = 0
        batch = db.batch()
        
        for note in notes_query:
            batch.delete(note.reference)
            deleted_count += 1
            # Firestore batches accept up to 500 ops. Simple app likely won't hit this limit instantly, 
            # but in production, you'd handle pagination. For now, this is safe.

        batch.commit() # Execute all note deletions

        # 3. Finally, delete the Subject itself
        sub_ref.delete()

        return {
            "status": "success", 
            "message": f"Deleted Subject '{subject_title}' and {deleted_count} associated notes."
        }

    except Exception as e:
        print(f"Delete Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)